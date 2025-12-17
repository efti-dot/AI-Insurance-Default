import os, base64, fitz, io
from pathlib import Path
from openai import OpenAI
from docx import Document
from pptx import Presentation
from dotenv import load_dotenv
from PIL import Image
import faiss
from utils import split_text, create_embeddings
import numpy as np

load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

class DocAI:
    def __init__(self, dim=1536):
        self.kb = []
        self.dim = dim
        self.index = faiss.IndexFlatL2(dim)
        self.metadata = []
        

    # ---------------- IMAGE HELPERS ----------------
    def img64(self, path):
        """Convert image file to base64 data URL"""
        ext = Path(path).suffix.lower().lstrip('.')
        mime = f"image/{'jpeg' if ext=='jpg' else ext}"
        data = base64.b64encode(open(path, 'rb').read()).decode()
        return f"data:{mime};base64,{data}"

    def img64_from_bytes(self, img_bytes):
        """Convert image bytes to base64 data URL"""
        try:
            img = Image.open(io.BytesIO(img_bytes))
            buffered = io.BytesIO()
            img.save(buffered, format="PNG")
            data = base64.b64encode(buffered.getvalue()).decode()
            return f"data:image/png;base64,{data}"
        except Exception as e:
            print(f"Error converting image: {e}")
            return None

    # ---------------- IMAGE ANALYSIS ----------------
    def analyze_img(self, path, name):
        """Send image to GPT-4o for description"""
        return client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": [
                {"type": "text", "text": f"Describe all details from {name}"},
                {"type": "image_url", "image_url": {"url": self.img64(path)}}
            ]}],
            max_tokens=1000
        ).choices[0].message.content

    def analyze_img_from_url(self, img_url, name, context=""):
        """Analyze image from base64 URL"""
        return client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": [
                {"type": "text", "text": f"Describe all details from image in {name}. {context}"},
                {"type": "image_url", "image_url": {"url": img_url}}
            ]}],
            max_tokens=1000
        ).choices[0].message.content

    # ---------------- DOCX ----------------
    def text_docx(self, p):
        """Extract text AND images from DOCX"""
        doc = Document(p)
        content_parts = []

        # Extract text
        text_content = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        if text_content:
            content_parts.append(f"TEXT CONTENT:\n{text_content}")

        # Extract images
        img_count = 0
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    img_bytes = rel.target_part.blob
                    img_url = self.img64_from_bytes(img_bytes)
                    if img_url:
                        img_count += 1
                        img_desc = self.analyze_img_from_url(
                            img_url,
                            p.name,
                            f"This is image {img_count} from the document"
                        )
                        content_parts.append(f"\nIMAGE {img_count} DESCRIPTION:\n{img_desc}")
                except Exception as e:
                    print(f"Error processing image: {e}")

        return "\n\n".join(content_parts) if content_parts else "No content found"

    # ---------------- PPTX ----------------
    def text_pptx(self, p):
        """Extract text AND images from PPTX"""
        prs = Presentation(p)
        content_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            # Extract text
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

            if slide_text:
                content_parts.append(f"SLIDE {slide_num} TEXT:\n" + "\n".join(slide_text))

            # Extract images
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        img_bytes = shape.image.blob
                        img_url = self.img64_from_bytes(img_bytes)
                        if img_url:
                            img_desc = self.analyze_img_from_url(
                                img_url,
                                p.name,
                                f"Image from slide {slide_num}"
                            )
                            content_parts.append(f"SLIDE {slide_num} IMAGE:\n{img_desc}")
                    except Exception as e:
                        print(f"Error processing image from slide {slide_num}: {e}")

        return "\n\n".join(content_parts) if content_parts else "No content found"

    # ---------------- PDF ----------------
    def text_pdf(self, p, name):
        """Extract text and images from PDF pages"""
        doc, out = fitz.open(p), []
        for i, pg in enumerate(doc):
            pix = pg.get_pixmap(matrix=fitz.Matrix(2, 2))
            b64 = base64.b64encode(pix.tobytes("png")).decode()
            r = client.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "user", "content": [
                    {"type": "text", "text": f"Extract all content (text and images) from page {i+1} of {name}"},
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}}
                ]}]
            )
            out.append(f"PAGE {i+1}:\n{r.choices[0].message.content}")
        return "\n\n".join(out)

    # ---------------- UPLOAD ----------------
    def upload(self, path):
        p = Path(path)
        ext = p.suffix.lower()

        print(f"Processing {p.name}...")

        if ext in [".png", ".jpg", ".jpeg"]:
            c = self.analyze_img(p, name=p.name)
            # embed description
            chunks = split_text(c)
            embeddings = create_embeddings(chunks)
            np_vectors = np.array([chunk["embedding"] for chunk in embeddings]).astype("float32")
            self.index.add(np_vectors)
            self.metadata.extend([chunk["text"] for chunk in embeddings])
            # also keep raw description in kb
            self.kb.append({"name": p.name, "content": c})
            return f"✓ {p.name} added (image description stored)"
        elif ext == ".pdf":
            c = self.text_pdf(p, p.name)
        elif ext == ".docx":
            c = self.text_docx(p)
        elif ext in [".pptx", ".ppt"]:
            c = self.text_pptx(p)
        else:
            return "Unsupported type"

        # embed text content
        chunks = split_text(c)
        embeddings = create_embeddings(chunks)
        np_vectors = np.array([chunk["embedding"] for chunk in embeddings]).astype("float32")
        self.index.add(np_vectors)
        self.metadata.extend([chunk["text"] for chunk in embeddings])

        self.kb.append({"name": p.name, "content": c})
        return f"✓ {p.name} added (content embedded)"


    # ---------------- ASK ----------------
    def ask(self, q):
        """Ask a question using all uploaded docs"""
        if not self.kb:
            return "❌ Upload docs first"

        ctx = "\n\n".join([f"=== {d['name']} ===\n{d['content']}" for d in self.kb])
        r = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "Answer using the document content provided. If images were described, reference those descriptions."},
                {"role": "user", "content": f"{ctx}\n\nQuestion: {q}"}
            ],
            temperature=0
        )
        return r.choices[0].message.content